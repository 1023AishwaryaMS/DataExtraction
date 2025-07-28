import os
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF


def extract_from_docx(file_path):
    doc = Document(file_path)
    text = '\n'.join([para.text for para in doc.paragraphs])
    return text


def extract_from_pptx(file_path):
    prs = Presentation(file_path)
    text_runs = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return '\n'.join(text_runs)


def extract_from_pdf(file_path):
    pdf = fitz.open(file_path)
    text = []
    for page_num in range(len(pdf)):
        page = pdf.load_page(page_num)
        text.append(page.get_text())
    pdf.close()
    return '\n'.join(text)


def extract_text(file_path):
    ext = file_path.lower().split('.')[-1]
    if ext == 'docx':
        return extract_from_docx(file_path)
    elif ext == 'pptx':
        return extract_from_pptx(file_path)
    elif ext == 'pdf':
        return extract_from_pdf(file_path)
    else:
        raise ValueError("Unsupported file type!")


def generate_video_script(text, max_words_per_segment=100):
    """
    Simple function to convert text into a video script
    that splits the text into chunks (segments)
    suitable for voiceover or narration.
    """
    words = text.split()
    segments = []
    for i in range(0, len(words), max_words_per_segment):
        segment = ' '.join(words[i:i + max_words_per_segment])
        segments.append(segment.strip())

    script = ""
    for idx, segment in enumerate(segments, start=1):
        script += f"Scene {idx}:\n"
        script += segment + "\n\n"
    return script


def main():
    print("Choose input method:")
    print("1: Upload and extract text from DOCX/PPTX/PDF")
    print("2: Paste direct text input")
    choice = input("Enter your choice (1 or 2): ").strip()

    if choice == '1':
        file_path = input("Enter full path to your DOCX/PPTX/PDF file: ").strip()
        if not os.path.exists(file_path):
            print("File not found!")
            return
        try:
            extracted_text = extract_text(file_path)
        except ValueError as e:
            print(e)
            return
    elif choice == '2':
        print("Paste your text input (end with a blank line):")
        lines = []
        while True:
            line = input()
            if line == '':
                break
            lines.append(line)
        extracted_text = '\n'.join(lines)
    else:
        print("Invalid choice.")
        return

    if not extracted_text.strip():
        print("No text extracted or provided.")
        return

    print("\nGenerating video script...\n")
    video_script = generate_video_script(extracted_text)

    output_filename = "video_script.txt"
    with open(output_filename, "w", encoding="utf-8") as f:
        f.write(video_script)

    print(f"Video script saved to {output_filename}")
    print("\n--- Video Script Preview ---\n")
    print(video_script)


if __name__ == "__main__":
    main()