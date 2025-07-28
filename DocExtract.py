import os
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
from PIL import Image

def extract_from_docx(file_path, img_output_folder):
    doc = Document(file_path)
    text = '\n'.join([para.text for para in doc.paragraphs])

    if not os.path.exists(img_output_folder):
        os.makedirs(img_output_folder)
    img_count = 0
    rels = doc.part.rels
    for rel in rels:
        rel = rels[rel]
        if "image" in rel.target_ref:
            img_count += 1
            img = rel.target_part.blob
            img_ext = rel.target_ref.split('.')[-1]
            img_filename = f'{os.path.basename(file_path)}_image_{img_count}.{img_ext}'
            img_path = os.path.join(img_output_folder, img_filename)
            with open(img_path, 'wb') as f:
                f.write(img)
            print(f"Saved DOCX image: {img_path}")
    return text

def extract_from_pptx(file_path, img_output_folder):
    prs = Presentation(file_path)
    text_runs = []
    if not os.path.exists(img_output_folder):
        os.makedirs(img_output_folder)
    img_count = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            print(f"Slide {slide_num} shape type: {shape.shape_type}")
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
            # Try extracting from any shape with an image attribute
            if hasattr(shape, "image"):
                img_count += 1
                image = shape.image
                img_bytes = image.blob
                img_ext = image.ext
                img_filename = f'{os.path.basename(file_path)}_slide{slide_num}_image_{img_count}.{img_ext}'
                img_path = os.path.join(img_output_folder, img_filename)
                with open(img_path, 'wb') as f:
                    f.write(img_bytes)
                print(f"Saved PPTX image: {img_path}")
    return '\n'.join(text_runs)

def extract_from_pdf(file_path, img_output_folder):
    pdf = fitz.open(file_path)
    text = []
    if not os.path.exists(img_output_folder):
        os.makedirs(img_output_folder)
    img_count = 0

    for page_num in range(len(pdf)):
        page = pdf.load_page(page_num)
        text.append(page.get_text())

        images = page.get_images(full=True)
        for img_index, img in enumerate(images, start=1):
            xref = img[0]
            base_image = pdf.extract_image(xref)
            img_bytes = base_image["image"]
            img_ext = base_image["ext"]
            img_count += 1
            img_filename = f'{os.path.basename(file_path)}_page{page_num+1}_image_{img_count}.{img_ext}'
            img_path = os.path.join(img_output_folder, img_filename)
            with open(img_path, 'wb') as f:
                f.write(img_bytes)
            print(f"Saved PDF image: {img_path}")

    pdf.close()
    return '\n'.join(text)

def extract_all_docs_in_folder(folder_path):
    for filename in os.listdir(folder_path):
        filepath = os.path.join(folder_path, filename)

        if os.path.isfile(filepath):
            ext = filename.lower().split('.')[-1]
            if ext == 'docx':
                print(f"\nExtracting from DOCX file: {filename}")
                text = extract_from_docx(filepath, "extracted_images")
                print(text)
            elif ext == 'pptx':
                print(f"\nExtracting from PPTX file: {filename}")
                text = extract_from_pptx(filepath, "extracted_images")
                print(text)
            elif ext == 'pdf':
                print(f"\nExtracting from PDF file: {filename}")
                text = extract_from_pdf(filepath, "extracted_images")
                print(text)
            else:
                print(f"\nSkipping unsupported file type: {filename}")
# ...existing code...
# ...existing code...
if __name__ == "__main__":
    # To process a single file:
    file_to_scan = r"C:\Users\23200\OneDrive\Desktop\Python\DocumentExtraction\folder_to_scan\Automated Bird Species Identification using Audio Signal Processing.pptx"
    ext = file_to_scan.lower().split('.')[-1]
    if ext == 'docx':
        text = extract_from_docx(file_to_scan, "extracted_images")
    elif ext == 'pptx':
        text = extract_from_pptx(file_to_scan, "extracted_images")
    elif ext == 'pdf':
        text = extract_from_pdf(file_to_scan, "extracted_images")
    else:
        print("Unsupported file type!")
        text = ""
    print(text)

    # To process all files in a folder:
    folder_to_scan = r"C:\Users\23200\OneDrive\Desktop\Python\DocumentExtraction\folder_to_scan"
    extract_all_docs_in_folder(folder_to_scan)